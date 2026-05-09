"""Generate House Price Prediction presentation (v2 — no-overlap pass)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ─────────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x1A, 0x23, 0x4E)
PURPLE    = RGBColor(0x66, 0x7E, 0xEA)
L_PURPLE  = RGBColor(0xA0, 0xAE, 0xF5)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
L_GRAY    = RGBColor(0xEF, 0xF1, 0xF7)
D_GRAY    = RGBColor(0x28, 0x28, 0x38)
MID_GRAY  = RGBColor(0x66, 0x66, 0x77)
GREEN     = RGBColor(0x27, 0xAE, 0x60)
ORANGE    = RGBColor(0xD4, 0x74, 0x1A)
RED_C     = RGBColor(0xBD, 0x35, 0x28)
TEAL      = RGBColor(0x0E, 0x88, 0x6E)
CARD_BG   = RGBColor(0x22, 0x2D, 0x62)
VIOLE     = RGBColor(0x7D, 0x3C, 0x98)
ICON_BLU  = RGBColor(0x1A, 0x4A, 0x70)
ICON_GRN  = RGBColor(0x1C, 0x5E, 0x35)


def R(slide, l, t, w, h, clr):
    s = slide.shapes.add_shape(
        1, Inches(l), Inches(t), Inches(w), Inches(h)
    )
    s.fill.solid()
    s.fill.fore_color.rgb = clr
    s.line.fill.background()
    return s


def T(slide, text, l, t, w, h,
      sz=13, bold=False, clr=WHITE,
      align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    rn = p.add_run()
    rn.text = str(text)
    rn.font.size = Pt(sz)
    rn.font.bold = bold
    rn.font.color.rgb = clr
    rn.font.italic = italic
    return tb


def header(slide, title, sub=None):
    R(slide, 0, 0, 13.33, 1.20, NAVY)
    R(slide, 0, 1.20, 13.33, 0.06, PURPLE)
    T(slide, title, 0.35, 0.08, 12.6, 0.72, sz=26, bold=True, clr=WHITE)
    if sub:
        T(slide, sub, 0.35, 0.74, 12.6, 0.40, sz=12, clr=L_PURPLE, italic=True)


# ============================================================
# Build presentation
# ============================================================
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.50)
BLANK = prs.slide_layouts[6]

# ─────────────────────────────────────────────────────────────
# SLIDE 1  Title
# ─────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
R(s, 0, 0, 13.33, 7.5, NAVY)
R(s, 0, 0,    13.33, 0.16, PURPLE)
R(s, 0, 7.34, 13.33, 0.16, PURPLE)
R(s, 2.8, 3.22, 7.73, 0.06, PURPLE)
R(s, 2.8, 4.60, 7.73, 0.06, PURPLE)

T(s, "HOUSE PRICE PREDICTION",
  0.5, 1.10, 12.33, 1.10,
  sz=46, bold=True, clr=WHITE, align=PP_ALIGN.CENTER)
T(s, "End-to-End Machine Learning System",
  0.5, 2.22, 12.33, 0.65,
  sz=22, clr=L_PURPLE, align=PP_ALIGN.CENTER)
T(s, "Data Engineering  \u00b7  Feature Engineering  \u00b7  ML Model Training",
  0.5, 3.34, 12.33, 0.50,
  sz=15, clr=L_PURPLE, align=PP_ALIGN.CENTER)
T(s, "FastAPI Backend  \u00b7  PostgreSQL  \u00b7  Live Streamlit Dashboard",
  0.5, 3.90, 12.33, 0.50,
  sz=15, clr=L_PURPLE, align=PP_ALIGN.CENTER)
T(s, "Group Project  |  May 2026",
  0.5, 5.75, 12.33, 0.45,
  sz=15, clr=RGBColor(0x88, 0x99, 0xBB), align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────
# SLIDE 2  Project Overview
# ─────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
R(s, 0, 0, 13.33, 7.5, L_GRAY)
header(s, "Project Overview", "What we built and why")

# left panel
R(s, 0.35, 1.38, 6.35, 5.86, WHITE)
R(s, 0.35, 1.46, 0.06, 1.50, PURPLE)
T(s, "Mission", 0.58, 1.44, 5.9, 0.40, sz=15, bold=True, clr=NAVY)
T(s,
  "Predict US house prices for any address using a production-quality "
  "ML pipeline backed by real US government data APIs \u2014 fully free, "
  "fully generalisable nationally.",
  0.58, 1.88, 5.9, 0.82, sz=12, clr=D_GRAY)

R(s, 0.35, 2.90, 0.06, 4.05, PURPLE)
T(s, "Project Scope", 0.58, 2.88, 5.9, 0.40, sz=15, bold=True, clr=NAVY)
scope = [
    "22,098 training rows \u2014 King County WA + Ames Iowa datasets",
    "16-feature optimised model (property + census + school signals)",
    "FastAPI REST backend with 7 Alembic database migrations",
    "Streamlit live dashboard deployed on Render.com",
    "30,379 national ZIP centroid neighbourhood scorer",
    "92 automated tests passing",
]
for i, line in enumerate(scope):
    T(s, f"  \u2022  {line}", 0.56, 3.36 + i * 0.46, 6.0, 0.40, sz=12, clr=D_GRAY)

# right stats 2x3 grid
R(s, 7.08, 1.38, 5.92, 5.86, NAVY)
T(s, "KEY METRICS", 7.08, 1.44, 5.92, 0.38,
  sz=13, bold=True, clr=PURPLE, align=PP_ALIGN.CENTER)

stats  = [("R\u00b2 Score","0.9266"), ("MAE","$16,574"),
          ("RMSE","$20,904"),         ("Features","16"),
          ("Tests","92"),             ("ZCTAs","30,379")]
snotes = ["92.66% test accuracy","Mean absolute error",
          "Root mean sq. error","Optimised model inputs",
          "Automated tests","National ZIP coverage"]

for i, ((lbl, val), note) in enumerate(zip(stats, snotes)):
    col = i % 2
    row = i // 2
    cx = 7.16 + col * 2.92
    cy = 1.94 + row * 1.72
    R(s, cx, cy, 2.76, 1.52, CARD_BG)
    T(s, lbl,  cx+0.14, cy+0.09, 2.50, 0.30, sz=10, bold=True, clr=L_PURPLE)
    T(s, val,  cx+0.14, cy+0.40, 2.50, 0.52,
      sz=23, bold=True, clr=WHITE, align=PP_ALIGN.CENTER)
    T(s, note, cx+0.08, cy+1.16, 2.60, 0.28,
      sz=9, italic=True, clr=RGBColor(0x88, 0x99, 0xBB),
      align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────
# SLIDE 3  Problem & Data
# ─────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
R(s, 0, 0, 13.33, 7.5, L_GRAY)
header(s, "The Problem & Our Data",
       "Why house price prediction is hard \u2014 and how we solved it")

R(s, 0.35, 1.38, 6.10, 5.86, WHITE)
T(s, "Challenges We Addressed", 0.55, 1.45, 5.80, 0.40,
  sz=14, bold=True, clr=NAVY)
challenges = [
    ("Fragmented data",    "No single dataset covers all US markets"),
    ("Paid API lock-in",   "Most geocoding & property APIs require paid keys"),
    ("Local overfitting",  "KC OHE 'Neighborhood' (74 cols) = zero signal nationally"),
    ("Neighbourhood data", "Hard to quantify local market context objectively"),
    ("API reliability",    "Geocoders time-out; Census API has latency spikes"),
    ("Schema drift",       "DB auto-create diverged from migration \u2192 500 errors"),
]
for i, (prob, detail) in enumerate(challenges):
    y = 1.96 + i * 0.72
    R(s, 0.46, y+0.07, 0.15, 0.36, RED_C)
    T(s, prob,   0.70, y+0.02, 2.60, 0.30, sz=12, bold=True, clr=D_GRAY)
    T(s, detail, 0.70, y+0.32, 5.60, 0.34, sz=10, italic=True, clr=MID_GRAY)

R(s, 6.82, 1.38, 6.18, 5.86, NAVY)
T(s, "DATA SOURCES", 6.82, 1.44, 6.18, 0.38,
  sz=13, bold=True, clr=PURPLE, align=PP_ALIGN.CENTER)
sources = [
    ("King County WA",  "21,598 rows \u2014 historic MLS sales 2014-2015"),
    ("Ames, Iowa",      "500 rows \u2014 academic benchmark dataset"),
    ("US Census ACS5",  "Tract income, home value, occupancy rate"),
    ("Nominatim OSM",   "Free geocoding (address \u2192 lat/lon)"),
    ("FCC API",         "Census tract lookup from coordinates"),
    ("ZCTA centroids",  "30,379 national ZIP centroids for KNN scorer"),
]
for i, (name, desc) in enumerate(sources):
    y = 1.96 + i * 0.86
    R(s, 6.92, y, 5.96, 0.72, CARD_BG)
    T(s, name, 7.08, y+0.08, 3.0, 0.28, sz=13, bold=True, clr=WHITE)
    T(s, desc, 7.08, y+0.38, 5.7, 0.28, sz=10, italic=True, clr=L_PURPLE)

# ─────────────────────────────────────────────────────────────
# SLIDE 4  Feature Engineering
# ─────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
R(s, 0, 0, 13.33, 7.5, L_GRAY)
header(s, "Feature Engineering", "From raw address \u2192 16 model-ready signals")

R(s, 0.35, 1.38, 8.55, 5.86, WHITE)
T(s, "Feature Importance \u2014 Top 10  (87.5% cumulative)",
  0.55, 1.46, 8.20, 0.40, sz=14, bold=True, clr=NAVY)

feats = [
    ("GrLivArea \u2014 Above-ground living area",       14.12, PURPLE),
    ("LotArea \u2014 Total lot size (sq ft)",            10.84, PURPLE),
    ("MedianIncome \u2014 Census ACS5 tract",            10.81, TEAL),
    ("UnemploymentRate \u2014 Census ACS5",              10.77, TEAL),
    ("SchoolDistrictRating \u2014 Derived score",         8.80, ORANGE),
    ("GarageArea \u2014 Garage square feet",              8.77, PURPLE),
    ("YearBuilt \u2014 Construction year",                6.75, PURPLE),
    ("YearRemodAdd \u2014 Last renovation",               6.65, PURPLE),
    ("OverallQual \u2014 Build quality (1-10)",            6.65, PURPLE),
    ("TotRmsAbvGrd \u2014 Total rooms above grade",        3.28, PURPLE),
]
LABEL_W = 3.36
BAR_ORG = 3.76
MAX_BAR = 3.88
for i, (label, pct, clr) in enumerate(feats):
    y = 1.96 + i * 0.48
    T(s, label, 0.50, y, LABEL_W, 0.38, sz=11, clr=D_GRAY)
    bw = (pct / 16.0) * MAX_BAR
    R(s, BAR_ORG, y+0.07, bw, 0.26, clr)
    T(s, f"{pct:.1f}%", BAR_ORG + bw + 0.10, y+0.04, 0.72, 0.32,
      sz=10, bold=True, clr=D_GRAY)

# right top: categories
R(s, 9.22, 1.38, 3.78, 2.74, NAVY)
T(s, "FEATURE CATEGORIES", 9.22, 1.44, 3.78, 0.38,
  sz=11, bold=True, clr=PURPLE, align=PP_ALIGN.CENTER)
cats = [
    (PURPLE, "Property / Structural", "13 feat  \u2014  54.1% importance"),
    (TEAL,   "Census Economic",       "2 feat   \u2014  21.6% importance"),
    (ORANGE, "School District",       "1 feat   \u2014   8.8% importance"),
]
for i, (clr, cat, note) in enumerate(cats):
    cy = 1.94 + i * 0.72
    R(s, 9.28, cy, 3.64, 0.12, clr)
    T(s, cat,  9.28, cy+0.16, 3.64, 0.28, sz=12, bold=True, clr=WHITE)
    T(s, note, 9.28, cy+0.44, 3.64, 0.24, sz=10, clr=L_PURPLE)

# right bottom: engineered features
R(s, 9.22, 4.22, 3.78, 3.02, WHITE)
T(s, "ENGINEERED FEATURES", 9.36, 4.28, 3.56, 0.38,
  sz=11, bold=True, clr=NAVY)
eng = [
    "Age = 2026 \u2212 YearBuilt",
    "TotalSF = GrLivArea + BasementSF",
    "TotalBath = FullBath + 0.5 \u00d7 HalfBath",
    "QualArea = OverallQual \u00d7 GrLivArea",
    "log(LotArea) \u2014 reduces right skew",
    "p99 price cap \u2014 removes ~216 outliers",
]
for i, ef in enumerate(eng):
    T(s, f"\u2022 {ef}", 9.36, 4.72 + i*0.40, 3.60, 0.36, sz=10, clr=D_GRAY)

# ─────────────────────────────────────────────────────────────
# SLIDE 5  Model Development
# ─────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
R(s, 0, 0, 13.33, 7.5, L_GRAY)
header(s, "Model Development & Selection",
       "Systematic progression from baseline to production-ready")

R(s, 0.35, 1.38, 8.02, 5.86, WHITE)
T(s, "Model Progression", 0.55, 1.45, 7.72, 0.40, sz=14, bold=True, clr=NAVY)

# table columns
CX = [0.42,  1.62,  4.22,  5.12,  6.08,  7.04]
CW = [1.16,  2.56,  0.86,  0.92,  0.92,  1.28]
R(s, 0.38, 1.96, 7.96, 0.40, PURPLE)
hdrs = ["Phase", "Key Change", "# Feats", "Test R\u00b2", "MAE", "Status"]
for cx, cw, hdr in zip(CX, CW, hdrs):
    T(s, hdr, cx, 2.00, cw, 0.32, sz=11, bold=True,
      clr=WHITE, align=PP_ALIGN.CENTER)

rows = [
    ("Baseline",   "Random Forest, 13 property features",    "13", "0.917",  "$43k",   "Initial"),
    ("Phase 1",    "+ School District Rating",                "14", "0.920",  "$38k",   "Better"),
    ("Phase 2",    "+ Census MedianIncome + UnemployRate",    "16", "0.9266", "$16.6k", "OPTIMAL"),
    ("Phase 3",    "+ RenterOccupied (tested only)",          "17", "0.9258", "$17k",   "Overfit"),
    ("Production", "LightGBM 645-est Huber, 22 features",    "22", "0.879",  "$62k",   "Live API"),
]
row_bgs   = [WHITE, WHITE,
             RGBColor(0xDF, 0xF5, 0xE7),
             RGBColor(0xFB, 0xE8, 0xE8),
             RGBColor(0xEB, 0xF0, 0xFF)]
stat_clrs = [D_GRAY, GREEN,
             RGBColor(0x0A, 0x7A, 0x3A),
             RED_C, TEAL]
for ri, (row, bg) in enumerate(zip(rows, row_bgs)):
    ry = 2.40 + ri * 0.88
    R(s, 0.38, ry, 7.96, 0.82, bg)
    for j, (val, cx, cw) in enumerate(zip(row, CX, CW)):
        fc  = stat_clrs[ri] if j == 5 else D_GRAY
        bld = (j == 5) or (ri == 2 and j in [3, 4])
        T(s, val, cx, ry+0.18, cw, 0.46,
          sz=11, clr=fc, bold=bld, align=PP_ALIGN.CENTER)

# right panels
R(s, 8.72, 1.38, 4.28, 2.68, NAVY)
T(s, "ALGORITHMS TRIED", 8.88, 1.44, 3.96, 0.38,
  sz=12, bold=True, clr=PURPLE, align=PP_ALIGN.CENTER)
algos = [
    ("Random Forest",    "Baseline \u2014 scale-invariant, interpretable"),
    ("LightGBM",         "Production \u2014 gradient boost, fast"),
    ("Ridge Regression", "Linear benchmark for non-linearity"),
    ("Feature Ablation", "Systematic importance & overfit checks"),
]
for i, (algo, note) in enumerate(algos):
    y = 1.96 + i * 0.54
    T(s, algo, 8.88, y,        3.96, 0.30, sz=12, bold=True, clr=WHITE)
    T(s, note, 8.88, y+0.28, 3.96, 0.22, sz=10, italic=True, clr=L_PURPLE)

R(s, 8.72, 4.20, 4.28, 3.04, WHITE)
T(s, "KEY DESIGN DECISIONS", 8.86, 4.26, 4.10, 0.40, sz=12, bold=True, clr=NAVY)
decisions = [
    "Removed KC 'Neighborhood' OHE (74 cols, zero national signal)",
    "p99 price cap removes ~216 ultra-luxury outliers",
    "Early stopping: 645 optimal estimators (was 1500)",
    "Huber objective (\u03b1=0.9) \u2014 robust to residual outliers",
    "SimpleImputer preserves census slots for live retrain",
]
for i, d in enumerate(decisions):
    T(s, f"\u2022 {d}", 8.86, 4.72 + i*0.46, 4.08, 0.40, sz=10, clr=D_GRAY)

# ─────────────────────────────────────────────────────────────
# SLIDE 6  Architecture
# ─────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
R(s, 0, 0, 13.33, 7.5, L_GRAY)
header(s, "System Architecture",
       "Request-to-prediction pipeline with full observability")

# pipeline boxes
BOXES = [
    ("Client\nRequest",   "Address input",            0.22, 2.90, 2.00, 1.06, PURPLE),
    ("Geocoding",         "Nominatim + FCC\nlat/lon + tract",  2.54, 2.66, 2.10, 1.30, TEAL),
    ("Feature\nAssembly", "Census ACS5\n+ KNN scorer",         4.98, 2.66, 2.18, 1.30, ORANGE),
    ("ML Inference",      "LightGBM\n22 features",             7.50, 2.66, 2.10, 1.30, VIOLE),
    ("API\nResponse",     "Price + audit\n+ policy meta",      9.96, 2.90, 3.10, 1.06, NAVY),
]
for (title, sub, bx, by, bw, bh, bc) in BOXES:
    R(s, bx, by, bw, bh, bc)
    T(s, title, bx+0.10, by+0.07, bw-0.18, 0.40,
      sz=13, bold=True, clr=WHITE, align=PP_ALIGN.CENTER)
    T(s, sub,   bx+0.10, by+0.52, bw-0.18, 0.70,
      sz=10, clr=WHITE, align=PP_ALIGN.CENTER)

arrow_centers = [2.30, 4.74, 7.26, 9.66]
for ax in arrow_centers:
    T(s, "\u25b6", ax+0.04, 3.08, 0.28, 0.36,
      sz=14, bold=True, clr=NAVY, align=PP_ALIGN.CENTER)

# DB layer
R(s, 0.35, 4.22, 12.62, 0.96, NAVY)
T(s,
  "PostgreSQL / SQLite  \u00b7  SQLAlchemy ORM  \u00b7  7 Alembic Migrations"
  "  \u00b7  Prediction Reuse Cache  \u00b7  Workflow Events Log",
  0.55, 4.40, 12.2, 0.60, sz=12, clr=WHITE, align=PP_ALIGN.CENTER)

# bottom panels (equal width)
R(s, 0.35, 5.36, 6.06, 1.90, WHITE)
T(s, "Feature Policy Engine", 0.52, 5.42, 5.86, 0.36, sz=13, bold=True, clr=NAVY)
T(s, "3 profiles: balanced-v1  \u00b7  quality-first-v1  \u00b7  land-first-v1",
  0.52, 5.82, 5.86, 0.28, sz=11, clr=D_GRAY)
T(s, "Per-request policy override  \u00b7  State-targeted routing  \u00b7  422 on unknown policy",
  0.52, 6.14, 5.86, 0.28, sz=11, clr=D_GRAY)
T(s, "Policy metadata surfaced in every prediction response + governance audit",
  0.52, 6.46, 5.86, 0.28, sz=11, clr=D_GRAY)

R(s, 6.72, 5.36, 6.26, 1.90, WHITE)
T(s, "Provider Resilience", 6.88, 5.42, 6.06, 0.36, sz=13, bold=True, clr=NAVY)
T(s, "Retry budget + NonRetryableProviderError short-circuit for fast fail",
  6.88, 5.82, 6.06, 0.28, sz=11, clr=D_GRAY)
T(s, "Default timeout 8 s  \u00b7  Heuristic fallback when Census unavailable",
  6.88, 6.14, 6.06, 0.28, sz=11, clr=D_GRAY)
T(s, "Nominatim city/state fallback on no-match results (was 502 error)",
  6.88, 6.46, 6.06, 0.28, sz=11, clr=D_GRAY)

# ─────────────────────────────────────────────────────────────
# SLIDE 7  FastAPI Backend
# ─────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
R(s, 0, 0, 13.33, 7.5, L_GRAY)
header(s, "FastAPI Backend",
       "Production-grade REST API with observability built-in")

R(s, 0.35, 1.36, 6.52, 5.90, WHITE)
T(s, "REST Endpoints", 0.55, 1.44, 6.22, 0.40,
  sz=14, bold=True, clr=NAVY)

endpoints = [
    (PURPLE, "POST", "/v1/predictions",               "Predict price from address"),
    (TEAL,   "GET",  "/v1/predictions",               "List recent predictions"),
    (TEAL,   "GET",  "/v1/predictions/{id}",          "Full detail + feature vector"),
    (TEAL,   "GET",  "/v1/predictions/{id}/events",   "Workflow audit trace"),
    (ORANGE, "GET",  "/v1/dashboard",                 "Live metrics snapshot"),
    (ORANGE, "GET",  "/v1/health",                    "Service + model health"),
    (VIOLE,  "GET",  "/v1/policies",                  "Feature policy catalog"),
    (VIOLE,  "POST", "/v1/policies/simulate",         "Compare policy outcomes"),
    (RED_C,  "POST", "/v1/validate/address-baseline", "Scenario validation checks"),
]
for i, (clr, meth, path, desc) in enumerate(endpoints):
    ey = 1.96 + i * 0.54
    R(s, 0.46, ey+0.01, 0.68, 0.34, clr)
    T(s, meth, 0.46, ey+0.04, 0.68, 0.28,
      sz=9, bold=True, clr=WHITE, align=PP_ALIGN.CENTER)
    T(s, path, 1.20, ey,       3.30, 0.32, sz=11, bold=True, clr=D_GRAY)
    T(s, desc, 1.20, ey+0.31, 5.40, 0.20, sz=9,  italic=True, clr=MID_GRAY)

# right top
R(s, 7.22, 1.36, 5.78, 2.74, NAVY)
T(s, "TECH STACK", 7.38, 1.42, 5.46, 0.40,
  sz=13, bold=True, clr=PURPLE, align=PP_ALIGN.CENTER)
tech = [
    ("Python 3.12",  "FastAPI + Pydantic v2 + SQLAlchemy 2"),
    ("LightGBM",     "645 estimators, Huber objective"),
    ("PostgreSQL",   "Production DB  \u00b7  SQLite for dev/test"),
    ("Alembic",      "7 versioned schema migrations"),
    ("Render.com",   "Cloud deployment target"),
]
for i, (name, sub) in enumerate(tech):
    y = 1.96 + i * 0.44
    T(s, name, 7.38, y,        2.6, 0.28, sz=12, bold=True, clr=WHITE)
    T(s, sub,  7.38, y+0.26, 5.54, 0.20, sz=9,  italic=True, clr=L_PURPLE)

# right bottom
R(s, 7.22, 4.24, 5.78, 3.02, WHITE)
T(s, "OBSERVABILITY LAYER", 7.38, 4.30, 5.46, 0.40,
  sz=13, bold=True, clr=NAVY)
obs = [
    ("Workflow Events",   "Every lifecycle stage persisted (DB + paginated API)"),
    ("Prediction Reuse",  "Cache + policy-aware invalidation"),
    ("Feature Policy",    "Per-request policy + governance metadata"),
    ("Audit Trace",       "Full lineage: geocode \u2192 enrich \u2192 predict"),
    ("Health Endpoint",   "Model version, DB state, provider availability"),
    ("Feature Overrides", "Per-request inject + bypass-reuse flag"),
]
for i, (name, sub) in enumerate(obs):
    y = 4.76 + i * 0.40
    T(s, name, 7.38, y,        2.46, 0.26, sz=11, bold=True, clr=D_GRAY)
    T(s, sub,  9.88, y,        3.06, 0.26, sz=9,  italic=True, clr=MID_GRAY)

# ─────────────────────────────────────────────────────────────
# SLIDE 8  Dashboard
# ─────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
R(s, 0, 0, 13.33, 7.5, L_GRAY)
header(s, "Live Dashboard & UX",
       "Streamlit front-end, publicly deployed on Render.com")

R(s, 0.35, 1.36, 6.08, 5.90, WHITE)
T(s, "Dashboard Features", 0.55, 1.44, 5.78, 0.40,
  sz=14, bold=True, clr=NAVY)
dash = [
    ("Address-to-Price",     "Enter any US address \u2014 get instant ML prediction"),
    ("Feature Breakdown",    "All 16 model inputs with values & importance bars"),
    ("Market Context",       "Census ACS5 & neighbourhood score overlays"),
    ("Price Distribution",   "Interactive Plotly histograms & scatter charts"),
    ("Policy Simulation",    "Compare balanced vs quality-first vs land-first"),
    ("Prediction History",   "Browse recent requests with workflow audit traces"),
    ("Health Monitor",       "Live provider status, model version, DB health"),
    ("Workflow Events Feed", "Real-time lifecycle event log per prediction"),
]
for i, (name, desc) in enumerate(dash):
    dy = 1.96 + i * 0.61
    R(s, 0.46, dy+0.11, 0.15, 0.30, PURPLE)
    T(s, name, 0.70, dy+0.02, 3.0, 0.28, sz=12, bold=True, clr=D_GRAY)
    T(s, desc, 0.70, dy+0.30, 5.6, 0.28, sz=10, italic=True, clr=MID_GRAY)

# right mockup
R(s, 6.78, 1.36, 6.22, 5.90, NAVY)
T(s, "DASHBOARD PREVIEW", 6.94, 1.42, 5.90, 0.38,
  sz=12, bold=True, clr=PURPLE, align=PP_ALIGN.CENTER)

def ui_card(sl, px, py, pw, ph, pc, title, val, sub=None):
    R(sl, px, py, pw, ph, pc)
    T(sl, title, px+0.12, py+0.08, pw-0.22, 0.26,
      sz=9, bold=True, clr=PURPLE)
    T(sl, val, px+0.12, py+0.36, pw-0.22, ph-0.50, sz=12, clr=WHITE)
    if sub:
        T(sl, sub, px+0.12, py+ph-0.28, pw-0.22, 0.22,
          sz=9, italic=True, clr=L_PURPLE)

ui_card(s, 6.88, 1.90, 6.00, 0.72, CARD_BG,
        "Enter Address", "123 Main St, Seattle WA 98101")
ui_card(s, 6.88, 2.74, 2.88, 1.10, ICON_GRN,
        "Predicted Price", "$524,000", "\u2191 +2.3% vs median")
ui_card(s, 9.88, 2.74, 3.00, 1.10, ICON_BLU,
        "Confidence Band", "$482k \u2013 $566k", "R\u00b2 = 0.9266")

R(s, 6.88, 3.96, 6.00, 1.24, CARD_BG)
T(s, "Feature Importance", 7.00, 4.02, 5.76, 0.26,
  sz=9, bold=True, clr=PURPLE)
T(s,
  "GrLivArea     \u2588\u2588\u2588\u2588\u2588\u2588\u2588\u2588  14.1%\n"
  "MedianIncome  \u2588\u2588\u2588\u2588\u2588\u2588\u2588  10.8%\n"
  "SchoolRating  \u2588\u2588\u2588\u2588\u2588\u2588   8.8%",
  7.00, 4.30, 5.76, 0.84, sz=11, clr=WHITE)

ui_card(s, 6.88, 5.30, 6.00, 0.60, CARD_BG,
        "Active Policy", "balanced-v1  \u00b7  NeighbourhoodScore: 72 / 100")

T(s, "Live: house-price-prediction-1-vrwx.onrender.com",
  6.90, 7.04, 5.96, 0.28,
  sz=11, bold=True, clr=PURPLE, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────
# SLIDE 9  Results
# ─────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
R(s, 0, 0, 13.33, 7.5, L_GRAY)
header(s, "Model Results",
       "Performance metrics across development iterations")

mcards = [
    ("R\u00b2 Score", "0.9266",  "92.66% variance explained"),
    ("MAE",           "$16,574", "Mean absolute error (test set)"),
    ("RMSE",          "$20,904", "Root mean squared error"),
    ("Gap",           "7.28%",   "Train/test \u2014 great generalisation"),
]
for i, (lbl, val, note) in enumerate(mcards):
    cx = 0.35 + i * 3.22
    R(s, cx, 1.36, 3.08, 1.64, WHITE)
    R(s, cx, 1.36, 3.08, 0.09, PURPLE)
    T(s, lbl,  cx+0.14, 1.50, 2.80, 0.36, sz=14, bold=True, clr=NAVY)
    T(s, val,  cx+0.14, 1.88, 2.80, 0.52,
      sz=27, bold=True, clr=PURPLE, align=PP_ALIGN.CENTER)
    T(s, note, cx+0.14, 2.52, 2.80, 0.42,
      sz=10, italic=True, clr=D_GRAY, align=PP_ALIGN.CENTER)

R(s, 0.35, 3.22, 12.62, 4.02, WHITE)
T(s, "Training Iterations \u2014 Performance Evolution",
  0.55, 3.30, 12.0, 0.38, sz=14, bold=True, clr=NAVY)

EC = [0.42, 1.76, 5.90, 7.00, 8.06, 9.04]
EW = [1.30, 4.10, 1.06, 1.02, 0.94, 3.82]
EH = ["Version", "Key Change", "MAE", "RMSE", "R\u00b2", "Notes"]
R(s, 0.38, 3.76, 12.6, 0.40, PURPLE)
for cx, cw, hdr in zip(EC, EW, EH):
    T(s, hdr, cx, 3.80, cw, 0.30,
      sz=11, bold=True, clr=WHITE, align=PP_ALIGN.CENTER)

evol = [
    ("v1 \u2014 May 4",   "22k rows (KC+Ames), 20 features, Random Forest",      "$72k",   "$157k", "0.840", "First production model"),
    ("v2 \u2014 May 6a",  "Added Waterfront, BasementSF, ViewScore (22 feats)",   "$67k",   "$131k", "0.875", "Feature expansion"),
    ("v3 \u2014 May 6b",  "Interaction features, Huber, early stopping (LGB)",    "$62k",   "$97k",  "0.879", "RMSE dropped 26%!"),
    ("v4 \u2014 CSV 16f", "Optimised 16-feature offline (Ames + KC)",             "$16.6k", "$20.9k","0.9266","Best accuracy achieved"),
]
for ri, row in enumerate(evol):
    ry = 4.20 + ri * 0.74
    rbg = RGBColor(0xDF, 0xF5, 0xE7) if ri == 3 else (WHITE if ri % 2 == 0 else L_GRAY)
    R(s, 0.38, ry, 12.6, 0.68, rbg)
    for val, cx, cw in zip(row, EC, EW):
        fc  = GREEN if ri == 3 else D_GRAY
        bld = ri == 3
        T(s, val, cx, ry+0.11, cw, 0.46,
          sz=10, clr=fc, bold=bld, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────
# SLIDE 10  Challenges & Solutions
# ─────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
R(s, 0, 0, 13.33, 7.5, L_GRAY)
header(s, "Challenges & Solutions",
       "What we really learned building for national scale")

items = [
    (RED_C,  "National Generalisation",
     "KC 'Neighborhood' OHE (74 columns) = zero signal for non-KC addresses \u2014 silent failure mode.",
     "Removed OHE. Replaced with Census ACS5 live API + KNN scorer seeded on 30,379 ZCTA centroids."),
    (ORANGE, "API Latency & Failures",
     "Census geocoder + ACS5 take 3-8 s. Nominatim returns no-match on partial addresses.",
     "Raised timeout to 8 s. NonRetryableProviderError short-circuits retry budget. City/state fallback added."),
    (PURPLE, "Prediction Consistency",
     "Create and detail endpoints returned different prices \u2014 rounding applied after DB persistence.",
     "Round price before persistence. Regression test asserts create and detail prices always match."),
    (TEAL,   "Overfitting Signal",
     "17-feature model degraded test R\u00b2 by 0.0008 adding RenterOccupied \u2014 textbook overfitting.",
     "Kept 16-feature model. LightGBM early stopping found 645 optimal estimators (was hand-tuned 1500)."),
    (NAVY,   "Schema Drift to Production",
     "DB auto-create at runtime diverged from Alembic migration state, causing 500 errors in production.",
     "Migration-first startup (create_schema=False default). Auto-create restricted to test env only."),
]
for i, (clr, title, problem, solution) in enumerate(items):
    cy = 1.38 + i * 1.18
    R(s, 0.35, cy, 12.62, 1.06, WHITE)
    R(s, 0.35, cy, 0.20,  1.06, clr)
    T(s, title, 0.68, cy+0.06, 6.10, 0.34, sz=13, bold=True, clr=NAVY)
    T(s, f"Problem: {problem}",
      0.68, cy+0.44, 6.10, 0.56, sz=10, italic=True, clr=RED_C)
    R(s, 7.18, cy+0.06, 5.72, 0.90, RGBColor(0xEB, 0xF7, 0xEB))
    T(s, "Solution:", 7.30, cy+0.06, 1.30, 0.28,
      sz=10, bold=True, clr=GREEN)
    T(s, solution, 7.30, cy+0.34, 5.52, 0.60, sz=10, clr=D_GRAY)

# ─────────────────────────────────────────────────────────────
# SLIDE 11  Testing & Quality
# ─────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
R(s, 0, 0, 13.33, 7.5, L_GRAY)
header(s, "Testing & Quality",
       "92 automated tests \u2014 from unit through end-to-end smoke")

R(s, 0.35, 1.36, 6.12, 5.90, WHITE)
T(s, "Test Coverage Areas", 0.55, 1.44, 5.82, 0.40,
  sz=14, bold=True, clr=NAVY)
test_cats = [
    (PURPLE, "API Contract Tests",     "POST predict, GET list/detail, 404, 422 validation"),
    (TEAL,   "Feature Policy Tests",   "Weighted scoring, policy bias, unknown-policy 422"),
    (ORANGE, "Workflow Event Tests",   "Pagination, sort order, event name filtering"),
    (GREEN,  "Prediction Reuse Tests", "Cache hit/miss, policy-boundary invalidation"),
    (NAVY,   "Smoke Tests",            "48-check production readiness (fake + live providers)"),
    (VIOLE,  "Regression Tests",       "Create/detail price consistency, rounding"),
    (RED_C,  "Feature Override Tests", "bypass_reuse flag, override keys in workflow events"),
    (TEAL,   "Schema Validation",      "Feature bounds, inverted bound rejection (422)"),
]
for i, (clr, name, desc) in enumerate(test_cats):
    ty = 1.96 + i * 0.60
    R(s, 0.46, ty+0.12, 0.15, 0.30, clr)
    T(s, name, 0.70, ty+0.02, 3.20, 0.28, sz=12, bold=True, clr=D_GRAY)
    T(s, desc, 0.70, ty+0.30, 5.60, 0.26, sz=9,  italic=True, clr=MID_GRAY)

# right top milestones
R(s, 6.82, 1.36, 6.18, 2.72, NAVY)
T(s, "MILESTONES", 6.98, 1.42, 5.86, 0.38,
  sz=13, bold=True, clr=PURPLE, align=PP_ALIGN.CENTER)
milestones = [
    ("Apr 13", "25", "Initial backend scaffold"),
    ("Apr 15", "83", "National scorer + orchestration"),
    ("Apr 27", "68", "Policy engine + overrides"),
    ("May  6", "92", "Production readiness pass"),
]
for i, (date, cnt, note) in enumerate(milestones):
    my = 1.94 + i * 0.54
    R(s, 6.90, my, 6.02, 0.46, CARD_BG)
    T(s, date, 7.02, my+0.08, 1.22, 0.28, sz=11, clr=L_PURPLE)
    T(s, cnt,  8.30, my+0.02, 1.40, 0.38,
      sz=20, bold=True, clr=WHITE, align=PP_ALIGN.CENTER)
    T(s, note, 9.80, my+0.08, 3.02, 0.28, sz=10, italic=True, clr=L_PURPLE)

# right bottom practices
R(s, 6.82, 4.22, 6.18, 3.04, WHITE)
T(s, "QUALITY PRACTICES", 6.98, 4.28, 5.86, 0.38, sz=13, bold=True, clr=NAVY)
practices = [
    "Pytest with SQLite in-memory DB for full test isolation",
    "FakePropertyDataClient for deterministic test oracle",
    "Smoke test: 48 checks across all production endpoints",
    "Regression: create-vs-detail price consistency assertion",
    "Training guardrail: TRAINING_MIN_ROWS prevents tiny models",
    "Makefile one-command: test + lint + train pipeline",
]
for i, p in enumerate(practices):
    T(s, f"\u2022 {p}", 6.98, 4.72 + i*0.38, 5.86, 0.34, sz=11, clr=D_GRAY)

# ─────────────────────────────────────────────────────────────
# SLIDE 12  Key Learnings
# ─────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
R(s, 0, 0, 13.33, 7.5, L_GRAY)
header(s, "Key Learnings", "What we would tell our past selves")

learnings = [
    (PURPLE, "1", "Feature Generalisation Beats Local Accuracy",
     "A Census ACS5-backed model generalising to all US addresses outperforms one fine-tuned for a single city. "
     "Remove local OHE columns early \u2014 they are silent zero-signal failure modes at national inference time."),
    (TEAL,   "2", "Free Public APIs Are Production-Grade",
     "Nominatim, US Census ACS5, and FCC APIs are reliable, rate-limit tolerant, and sufficient for national "
     "coverage. No paid commercial geocoder was needed at our scale."),
    (ORANGE, "3", "Schema Discipline Saves Hours of Debugging",
     "Running Alembic migration-first from day one would have prevented the production 500-error incident. "
     "Always treat DB schema as code \u2014 never a runtime convenience."),
    (GREEN,  "4", "Interaction Features & Outlier Capping Are High-ROI",
     "QualArea (OverallQual \u00d7 GrLivArea) + p99 price capping dropped RMSE by 26% in one pass \u2014 "
     "larger gain than adding an entire new data source."),
    (NAVY,   "5", "Observability Is a Feature, Not an Afterthought",
     "Workflow events, audit traces, policy metadata, and the health endpoint transformed a black-box API "
     "into a debuggable, governable system \u2014 essential for a multi-developer group project."),
]
for i, (clr, num, title, body) in enumerate(learnings):
    ly = 1.38 + i * 1.17
    R(s, 0.35, ly, 12.62, 1.06, WHITE)
    R(s, 0.35, ly, 0.54,  1.06, clr)
    T(s, num,   0.35, ly+0.22, 0.54, 0.54,
      sz=26, bold=True, clr=WHITE, align=PP_ALIGN.CENTER)
    T(s, title, 1.04, ly+0.06, 7.60, 0.36, sz=14, bold=True, clr=NAVY)
    T(s, body,  1.04, ly+0.46, 11.5, 0.54, sz=11, clr=D_GRAY)

# ─────────────────────────────────────────────────────────────
# SLIDE 13  Thank You
# ─────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
R(s, 0, 0, 13.33, 7.50, NAVY)
R(s, 0, 0,    13.33, 0.16, PURPLE)
R(s, 0, 7.34, 13.33, 0.16, PURPLE)

T(s, "Thank You",
  1.0, 1.85, 11.33, 1.20,
  sz=56, bold=True, clr=WHITE, align=PP_ALIGN.CENTER)
T(s, "Questions & Discussion",
  1.0, 3.10, 11.33, 0.66,
  sz=24, clr=L_PURPLE, align=PP_ALIGN.CENTER)
R(s, 2.8, 3.92, 7.73, 0.07, PURPLE)

summary_lines = [
    "R\u00b2 = 0.9266   \u00b7   MAE = $16,574   \u00b7   RMSE = $20,904",
    "22,098 training rows   \u00b7   16 optimised features   \u00b7   92 tests passing",
    "FastAPI  +  LightGBM  +  PostgreSQL  +  Streamlit  \u2014  Deployed on Render",
]
for i, line in enumerate(summary_lines):
    T(s, line, 1.0, 4.12 + i*0.46, 11.33, 0.40,
      sz=14, clr=L_PURPLE, align=PP_ALIGN.CENTER)

T(s, "house-price-prediction-1-vrwx.onrender.com",
  1.8, 5.78, 9.73, 0.50,
  sz=16, bold=True, clr=PURPLE, align=PP_ALIGN.CENTER)

# ── save ─────────────────────────────────────────────────────
OUT = "/workspaces/House_Price_Prediction/House_Price_Prediction_Presentation.pptx"
prs.save(OUT)
print(f"Saved: {OUT}  ({len(prs.slides)} slides)")
